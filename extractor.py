"""
extractor.py — Extract content from PDF / DOCX / PPTX / TXT.

Returns three separate lists for each document:
  • text_chunks  : List[str]  — raw text, split into word-based chunks
  • tables       : List[str]  — one string per table (raw, not summarised)
  • images       : List[str]  — one base64 JPEG string per image

Text is embedded directly.
Tables and images are first summarised by Gemini (in summariser.py),
then the summaries are embedded but the raw originals are stored for retrieval.

PDF extraction intentionally avoids unstructured/pdfminer because that stack can
break with errors like:
    cannot import name 'PSSyntaxError' from 'pdfminer.pdfparser'
Instead, PDFs are processed with PyMuPDF (fitz), which is more robust here.
"""

import base64
import io
from pathlib import Path
from typing import List, Tuple

import nltk

for _pkg in ["averaged_perceptron_tagger_eng", "punkt", "punkt_tab"]:
    nltk.download(_pkg, quiet=True)

from config import CHUNK_OVERLAP, CHUNK_SIZE


# ─────────────────────────────────────────────────────────────────────────────
# Helper: word-based sliding-window chunker
# ─────────────────────────────────────────────────────────────────────────────

def _chunk(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    words = text.split()
    chunks, i = [], 0
    step = max(1, size - overlap)
    while i < len(words):
        chunk = " ".join(words[i : i + size])
        if chunk.strip():
            chunks.append(chunk.strip())
        i += step
    return chunks


# ─────────────────────────────────────────────────────────────────────────────
# Helper: PIL image → base64 JPEG string
# ─────────────────────────────────────────────────────────────────────────────

def _pil_to_b64(img) -> str:
    buf = io.BytesIO()
    img.convert("RGB").save(buf, format="JPEG")
    return base64.b64encode(buf.getvalue()).decode("utf-8")


# ─────────────────────────────────────────────────────────────────────────────
# Per-format extractors
# ─────────────────────────────────────────────────────────────────────────────

def _from_pdf(file_path: Path) -> Tuple[List[str], List[str], List[str]]:
    """
    Extract PDF text and images with PyMuPDF.

    Why this implementation:
    - avoids the unstructured/pdfminer import path that is currently failing
    - works on regular PDFs without the PSSyntaxError issue

    Current trade-off:
    - text extraction is robust
    - embedded images are extracted
    - table extraction is not attempted here, so `tables` is returned empty for PDFs
      instead of crashing the whole pipeline
    """
    try:
        import fitz  # PyMuPDF
        from PIL import Image as PILImage
    except Exception as exc:
        raise ImportError(
            "PyMuPDF is required for PDF extraction. Install it with: pip install pymupdf"
        ) from exc

    raw_text_parts: List[str] = []
    tables: List[str] = []
    images: List[str] = []

    with fitz.open(file_path) as doc:
        for page_num, page in enumerate(doc, 1):
            text = page.get_text("text") or ""
            if text.strip():
                raw_text_parts.append(f"[Page {page_num}] {text.strip()}")

            for image_info in page.get_images(full=True):
                xref = image_info[0]
                try:
                    img_dict = doc.extract_image(xref)
                    img_bytes = img_dict.get("image")
                    if not img_bytes:
                        continue
                    img = PILImage.open(io.BytesIO(img_bytes))
                    images.append(_pil_to_b64(img))
                except Exception:
                    continue

    text_chunks = _chunk(" ".join(raw_text_parts))
    return text_chunks, tables, images


def _from_docx(file_path: Path) -> Tuple[List[str], List[str], List[str]]:
    from docx import Document
    from PIL import Image as PILImage

    doc = Document(file_path)
    raw_text_parts, tables, images = [], [], []

    for para in doc.paragraphs:
        if para.text.strip():
            raw_text_parts.append(para.text.strip())

    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            tables.append("\n".join(rows))

    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img = PILImage.open(io.BytesIO(rel.target_part.blob))
                images.append(_pil_to_b64(img))
            except Exception:
                pass

    return _chunk(" ".join(raw_text_parts)), tables, images


def _from_pptx(file_path: Path) -> Tuple[List[str], List[str], List[str]]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image as PILImage

    prs = Presentation(file_path)
    raw_text_parts, tables, images = [], [], []

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                raw_text_parts.append(f"[Slide {slide_num}] {shape.text.strip()}")

            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        rows.append(" | ".join(cells))
                if rows:
                    tables.append(f"[Slide {slide_num} Table]\n" + "\n".join(rows))

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = PILImage.open(io.BytesIO(shape.image.blob))
                    images.append(_pil_to_b64(img))
                except Exception:
                    pass

    return _chunk(" ".join(raw_text_parts)), tables, images


def _from_txt(file_path: Path) -> Tuple[List[str], List[str], List[str]]:
    text = file_path.read_text(encoding="utf-8", errors="ignore")
    return _chunk(text), [], []


# ─────────────────────────────────────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────────────────────────────────────

def extract(file_path: Path) -> Tuple[List[str], List[str], List[str]]:
    """
    Extract all content from a supported document.

    Returns
    -------
    text_chunks : List[str]  — plain text passages, already chunked
    tables      : List[str]  — full table strings, one per table
    images      : List[str]  — base64-encoded JPEG strings, one per image
    """
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        return _from_pdf(file_path)
    if ext in {".docx", ".doc"}:
        return _from_docx(file_path)
    if ext in {".pptx", ".ppt"}:
        return _from_pptx(file_path)
    if ext == ".txt":
        return _from_txt(file_path)
    raise ValueError(f"Unsupported format: {ext}")
